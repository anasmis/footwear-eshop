# -*- coding: utf-8 -*-
"""Génère la présentation PowerPoint du projet E-Shop Django.
Source : rapport.tex  |  Captures : screenshots/
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ---- Charte graphique (reprise du rapport) ----------------------------------
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
PINK   = RGBColor(0xDB, 0x27, 0x77)
DARK   = RGBColor(0x2D, 0x2D, 0x37)
GRAY   = RGBColor(0x6B, 0x6B, 0x78)
LIGHT  = RGBColor(0xF5, 0xF5, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def setrun(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def header(s, title, kicker=None):
    """Bandeau de titre standard pour une slide de contenu."""
    rect(s, 0, 0, SW, Inches(1.15), VIOLET)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), PINK)
    tf = textbox(s, Inches(0.6), Inches(0.12), Inches(12.1), Inches(0.95),
                 anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        p = tf.paragraphs[0]
        setrun(p, kicker, 12, RGBColor(0xE9, 0xD5, 0xFF), bold=True)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    setrun(p2, title, 28, WHITE, bold=True)


def footer(s, page):
    tf = textbox(s, Inches(0.6), Inches(7.02), Inches(8), Inches(0.4))
    setrun(tf.paragraphs[0], "E-Shop Django  •  Projet de fin de module 2025–2026",
           9, GRAY, italic=True)
    tf2 = textbox(s, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.4))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    setrun(tf2.paragraphs[0], str(page), 9, GRAY)


def bullets(s, items, x, y, w, h, size=16, gap=10, color=DARK):
    tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = lvl
        bullet = "▸ " if lvl == 0 else "– "
        setrun(p, bullet, size, PINK if lvl == 0 else VIOLET, bold=True)
        # support bold prefix via "**...**"
        if "||" in it:
            head, rest = it.split("||", 1)
            setrun(p, head, size, color, bold=True)
            setrun(p, rest, size, color)
        else:
            setrun(p, it, size, color)
    return tf


def add_image_fit(s, path, x, y, max_w, max_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w
    h = int(w / ar)
    if h > max_h:
        h = max_h
        w = int(h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    # cadre
    rect(s, px - Emu(9525), py - Emu(9525), w + Emu(19050), h + Emu(19050),
         WHITE, line=VIOLET)
    s.shapes.add_picture(path, px, py, width=w, height=h)


SS = "screenshots"
PAGE = [0]


def pg():
    PAGE[0] += 1
    return PAGE[0]


# =============================================================== 1. TITRE
s = slide()
rect(s, 0, 0, SW, SH, VIOLET)
rect(s, 0, Inches(2.55), SW, Inches(2.65), RGBColor(0x6D, 0x28, 0xD9))
rect(s, 0, Inches(5.0), SW, Inches(0.10), PINK)

tf = textbox(s, Inches(1), Inches(0.7), Inches(11.3), Inches(0.9))
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
setrun(tf.paragraphs[0], "Génie MIS  —  Module : Développement Web avec Django",
       16, RGBColor(0xE9, 0xD5, 0xFF))

tf = textbox(s, Inches(1), Inches(2.75), Inches(11.3), Inches(1.3),
             anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
setrun(tf.paragraphs[0], "E-Shop Django", 60, WHITE, bold=True)

tf = textbox(s, Inches(1), Inches(4.05), Inches(11.3), Inches(0.9),
             anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
setrun(tf.paragraphs[0],
       "Développement, sécurisation et déploiement d'une plateforme e-commerce moderne",
       20, RGBColor(0xF3, 0xE8, 0xFF), italic=True)

tf = textbox(s, Inches(1), Inches(5.45), Inches(11.3), Inches(1.4),
             anchor=MSO_ANCHOR.TOP)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
setrun(p, "Réalisé par  ", 18, WHITE)
setrun(p, "Ali Bouhamidi  &  Said Aghrod", 18, WHITE, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
setrun(p, "Encadré par  ", 16, RGBColor(0xE9, 0xD5, 0xFF))
setrun(p, "Prof. Bousselham", 16, RGBColor(0xE9, 0xD5, 0xFF), bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(10)
setrun(p, "Année universitaire 2025 – 2026", 14, RGBColor(0xC4, 0xB5, 0xFD))

# =============================================================== 2. PLAN
s = slide(); header(s, "Plan de la présentation")
left = ["01||  Contexte & problématique", "02||  Objectifs du projet",
        "03||  Analyse des besoins", "04||  Conception & modèle de données",
        "05||  Architecture & choix techniques"]
right = ["06||  Fonctionnalités de la plateforme", "07||  Intelligence artificielle (IA)",
         "08||  Sécurité", "09||  Déploiement & CI/CD",
         "10||  Démonstration & conclusion"]
bullets(s, left, Inches(0.8), Inches(1.7), Inches(5.9), Inches(5), size=18, gap=16)
bullets(s, right, Inches(6.9), Inches(1.7), Inches(5.9), Inches(5), size=18, gap=16)
footer(s, pg())

# =============================================================== 3. INTRO
s = slide(); header(s, "Introduction générale", "01  •  Contexte")
bullets(s, [
    "Le commerce électronique occupe une place centrale dans l'économie mondiale.",
    "E-Shop Django||  : une plateforme e-commerce complète développée comme projet de fin de module.",
    "Domaine retenu : une ||boutique de chaussures (dataset réel Myntra — fashion.csv).",
    "Catalogue organisé par ||type de chaussure (sport, ville, talons, ballerines…) et par rayons Homme / Femme.",
    "Objectif : une vraie ||application web professionnelle, pas un simple site vitrine.",
], Inches(0.8), Inches(1.7), Inches(11.7), Inches(5), size=18, gap=18)
footer(s, pg())

# =============================================================== 4. PROBLEMATIQUE
s = slide(); header(s, "Contexte & problématique", "01  •  Problématique")
box = rect(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.7), LIGHT, line=VIOLET)
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
setrun(tf.paragraphs[0], "Problématique", 15, VIOLET, bold=True)
p = tf.add_paragraph()
setrun(p, "Comment concevoir, développer, sécuriser et déployer une plateforme "
          "e-commerce moderne avec Django, tout en intégrant une fonctionnalité "
          "intelligente améliorant l'expérience utilisateur ?", 16, DARK, italic=True)
bullets(s, [
    "Organiser une application Django en modules cohérents",
    "Gérer produits, catégories, clients et commandes",
    "Sécuriser les accès selon les rôles des utilisateurs",
    "Gérer un panier et un processus de commande fiable",
    "Intégrer une fonctionnalité d'IA utile et préparer un déploiement réel",
], Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.3), size=16, gap=12)
footer(s, pg())

# =============================================================== 5. OBJECTIFS
s = slide(); header(s, "Objectifs du projet", "02  •  Objectifs")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.6))
setrun(tf.paragraphs[0], "Objectif général : ", 16, VIOLET, bold=True)
setrun(tf.paragraphs[0], "développer une plateforme e-commerce permettant de consulter "
       "des produits, gérer un panier, passer commande et bénéficier d'une recommandation IA.",
       16, DARK)
col1 = ["Architecture Django modulaire", "Modèles & base PostgreSQL",
        "Interfaces Templates + Bootstrap 5", "Inscription, connexion, profils",
        "Catalogue de produits", "Panier d'achat"]
col2 = ["Commandes & suivi de statut", "Tableau de bord administrateur",
        "Fonctionnalité d'IA (recommandation)", "Sécurisation de l'application",
        "Conteneurisation Docker / Compose", "Pipeline CI/CD (GitHub Actions)"]
bullets(s, col1, Inches(0.8), Inches(2.5), Inches(5.8), Inches(4.3), size=16, gap=12)
bullets(s, col2, Inches(6.9), Inches(2.5), Inches(5.8), Inches(4.3), size=16, gap=12)
footer(s, pg())

# =============================================================== 6. ACTEURS
s = slide(); header(s, "Analyse des besoins — les acteurs", "03  •  Besoins")
cards = [
    ("Visiteur", "Utilisateur non authentifié",
     ["Consulte accueil & catalogue", "Recherche / filtre les produits",
      "Voit le détail produit", "Crée un compte client"]),
    ("Client", "Utilisateur authentifié",
     ["Gère son profil", "Ajoute au panier", "Valide une commande",
      "Suit ses commandes & laisse un avis"]),
    ("Administrateur", "Gestionnaire de la plateforme",
     ["Gère produits, catégories, stock", "Fait évoluer le statut des commandes",
      "Consulte la liste des clients", "Accède aux statistiques"]),
]
cw = Inches(3.85); gap = Inches(0.3); x0 = Inches(0.7); y0 = Inches(1.6)
for i, (title, sub, its) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, Inches(0.95), VIOLET)
    rect(s, x, y0 + Inches(0.95), cw, Inches(4.35), LIGHT, line=VIOLET)
    tf = textbox(s, x, y0 + Inches(0.12), cw, Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    setrun(tf.paragraphs[0], title, 20, WHITE, bold=True)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    setrun(p, sub, 11, RGBColor(0xE9, 0xD5, 0xFF), italic=True)
    bullets(s, its, x + Inches(0.25), y0 + Inches(1.2), cw - Inches(0.5),
            Inches(3.9), size=13, gap=10)
footer(s, pg())

# =============================================================== 7. BESOINS FONCTIONNELS
s = slide(); header(s, "Besoins fonctionnels & non fonctionnels", "03  •  Besoins")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(5.8), Inches(0.5))
setrun(tf.paragraphs[0], "Fonctionnels", 18, VIOLET, bold=True)
bullets(s, [
    "Gestion produits & catégories",
    "Catalogue : recherche, filtres, tri, pagination",
    "Panier : ajout, quantité, total",
    "Commandes : checkout, statut, historique",
    "Avis clients (note + commentaire)",
    "Tableau de bord & recommandation IA",
], Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5), size=15, gap=11)
tf = textbox(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(0.5))
setrun(tf.paragraphs[0], "Non fonctionnels", 18, PINK, bold=True)
bullets(s, [
    "Sécurité (auth, CSRF, rôles)",
    "Performance (requêtes optimisées, cache IA)",
    "Portabilité (Docker)",
    "Maintenabilité (modulaire, tests)",
    "Ergonomie (responsive Bootstrap 5)",
], Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.5), size=15, gap=11, color=DARK)
footer(s, pg())

# =============================================================== 8. ARCHITECTURE
s = slide(); header(s, "Architecture applicative", "04  •  Conception")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.7))
setrun(tf.paragraphs[0], "Architecture Django modulaire (patron ", 16, DARK)
setrun(tf.paragraphs[0], "MVT", 16, VIOLET, bold=True)
setrun(tf.paragraphs[0], ") — chaque domaine métier isolé dans une application dédiée.",
       16, DARK)
apps = [
    ("ecommerce_project", "settings, urls, wsgi"),
    ("accounts", "utilisateurs, profils"),
    ("products", "produits, catégories, avis"),
    ("cart", "panier d'achat"),
    ("orders", "commandes & checkout"),
    ("dashboard", "statistiques & admin"),
    ("recommendation", "IA — TF-IDF + cosinus"),
    ("templates / static / media", "présentation & fichiers"),
]
cw = Inches(2.9); ch = Inches(1.25); gx = Inches(0.18); gy = Inches(0.25)
x0 = Inches(0.8); y0 = Inches(2.35)
for i, (name, desc) in enumerate(apps):
    r = i // 4; c = i % 4
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, LIGHT, line=VIOLET)
    tf = textbox(s, x + Inches(0.12), y + Inches(0.1), cw - Inches(0.24),
                 ch - Inches(0.2), anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], name, 14, VIOLET, bold=True, font="Consolas")
    p = tf.add_paragraph()
    setrun(p, desc, 11, GRAY)
footer(s, pg())

# =============================================================== 9. CHOIX TECHNO
s = slide(); header(s, "Choix technologiques", "04  •  Stack")
rows = [
    ("Backend", "Django 5.0 (Python 3.11)"),
    ("Frontend", "Templates Django + Bootstrap 5"),
    ("Base de données", "PostgreSQL 16 (SQLite en local)"),
    ("Authentification", "Système natif de Django"),
    ("Fichiers statiques", "WhiteNoise  •  médias : Pillow"),
    ("Intelligence artificielle", "scikit-learn — TF-IDF + similarité cosinus"),
    ("Serveur d'application", "Gunicorn (3 workers WSGI)"),
    ("Reverse proxy", "Nginx"),
    ("Conteneurisation", "Docker & Docker Compose"),
    ("CI/CD  &  Registre", "GitHub Actions  →  GHCR"),
]
y0 = Inches(1.5); rh = Inches(0.52); x0 = Inches(1.0)
w1 = Inches(4.2); w2 = Inches(7.1)
for i, (a, b) in enumerate(rows):
    y = y0 + i * rh
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, x0, y, w1, rh, VIOLET if False else bg, line=RGBColor(0xDD, 0xDD, 0xE5))
    rect(s, x0 + w1, y, w2, rh, bg, line=RGBColor(0xDD, 0xDD, 0xE5))
    tf = textbox(s, x0 + Inches(0.15), y, w1 - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], a, 13, VIOLET, bold=True)
    tf = textbox(s, x0 + w1 + Inches(0.15), y, w2 - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], b, 13, DARK)
footer(s, pg())

# =============================================================== 10. MODELE DE DONNEES
s = slide(); header(s, "Modèle de données", "04  •  Conception")
ents = [
    ("User / Profile", "Compte client + infos (téléphone, adresse)"),
    ("Category", "Catégorie / type de chaussure"),
    ("Product", "Nom, genre, prix, image, stock, dispo"),
    ("Cart / CartItem", "Panier persistant (1 par client) + lignes"),
    ("Order / OrderItem", "Commande validée + lignes (prix figés)"),
    ("Review", "Note 1–5 + commentaire (unique client/produit)"),
]
y0 = Inches(1.65); rh = Inches(0.78); x0 = Inches(0.9)
for i, (a, b) in enumerate(ents):
    y = y0 + i * rh
    rect(s, x0, y, Inches(0.18), rh - Inches(0.12), PINK)
    rect(s, x0 + Inches(0.28), y, Inches(11.3), rh - Inches(0.12), LIGHT,
         line=RGBColor(0xDD, 0xDD, 0xE5))
    tf = textbox(s, x0 + Inches(0.5), y, Inches(11), rh - Inches(0.12),
                 anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], a + "   —   ", 15, VIOLET, bold=True, font="Consolas")
    setrun(tf.paragraphs[0], b, 14, DARK)
footer(s, pg())

# =============================================================== 11. FONCTIONNALITES 1
s = slide(); header(s, "Fonctionnalités — comptes & catalogue", "06  •  Fonctionnalités")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5))
setrun(tf.paragraphs[0], "Comptes & rôles", 18, VIOLET, bold=True)
bullets(s, [
    "Inscription : compte créé + connexion auto, Profile généré par signal post_save",
    "Trois rôles : visiteur, client (authentifié), administrateur (is_staff)",
], Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4), size=15, gap=10)
tf = textbox(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.5))
setrun(tf.paragraphs[0], "Catalogue, recherche & filtres", 18, PINK, bold=True)
bullets(s, [
    "Recherche plein texte (nom + description)",
    "Filtres par catégorie et par genre (Homme / Femme)",
    "Tri (prix, date, nom) et pagination (9 produits / page)",
    "Seuls les produits disponibles sont affichés",
], Inches(0.8), Inches(3.95), Inches(11.7), Inches(2.8), size=15, gap=10)
footer(s, pg())

# =============================================================== 12. FONCTIONNALITES 2
s = slide(); header(s, "Fonctionnalités — panier & commandes", "06  •  Fonctionnalités")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(5.8), Inches(0.5))
setrun(tf.paragraphs[0], "Panier d'achat", 18, VIOLET, bold=True)
bullets(s, [
    "Panier persistant (1 Cart / utilisateur)",
    "Ajout, mise à jour, suppression, vidage",
    "Calcul du total",
    "Quantités bornées par le stock disponible",
], Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5), size=15, gap=11)
tf = textbox(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(0.5))
setrun(tf.paragraphs[0], "Commandes & checkout", 18, PINK, bold=True)
bullets(s, [
    "Checkout en transaction atomique",
    "Vérif. stock → création Order/OrderItem",
    "Prix & nom figés, décrément du stock",
    "6 statuts : attente → livrée / annulée",
    "Historique & détail des commandes",
], Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.5), size=15, gap=11, color=DARK)
footer(s, pg())

# =============================================================== 13. IA
s = slide(); header(s, "Intelligence artificielle — recommandation", "07  •  IA")
box = rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0), LIGHT, line=VIOLET)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3)
setrun(tf.paragraphs[0], "Recommandation de produits similaires", 15, VIOLET, bold=True)
p = tf.add_paragraph()
setrun(p, "Filtrage par contenu (content-based) affiché sous « Vous aimerez aussi » "
          "sur chaque fiche produit.", 14, DARK)
bullets(s, [
    "Chaque produit → document texte (nom pondéré + catégorie + description)",
    "TF-IDF|| entraîné sur 20 % du catalogue (n-grammes 1–2, stop-words FR)",
    "Catalogue complet projeté dans l'espace TF-IDF (transform)",
    "Similarité cosinus|| → les k produits les plus proches",
    "Mise en cache|| invalidée par la signature du catalogue + repli par catégorie",
], Inches(0.8), Inches(2.8), Inches(11.7), Inches(4), size=15, gap=12)
footer(s, pg())

# =============================================================== 14. SECURITE
s = slide(); header(s, "Sécurité", "08  •  Sécurité")
bullets(s, [
    "Authentification obligatoire|| (panier, checkout, avis) via @login_required",
    "Mots de passe hachés|| par Django + 4 password validators",
    "Protection CSRF|| active & validation des formulaires côté serveur",
    "Contrôle des rôles|| : dashboard/admin protégés par is_staff",
    "Variables sensibles|| (SECRET_KEY, BDD) dans un .env non versionné",
    "DEBUG=False en production|| : cookies sécurisés, HSTS, X-Frame-Options, nosniff",
    "Taille des uploads limitée|| (2,5 Mo Django / 5 Mo Nginx)",
], Inches(0.8), Inches(1.7), Inches(11.7), Inches(5), size=15, gap=13)
footer(s, pg())

# =============================================================== 15. DEPLOIEMENT
s = slide(); header(s, "Déploiement — Docker & Docker Compose", "09  •  Déploiement")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5))
setrun(tf.paragraphs[0], "Conteneurisation", 17, VIOLET, bold=True)
bullets(s, [
    "Image Docker basée sur python:3.11-slim",
    "entrypoint.sh : attend PostgreSQL → migrations → chargement des données → collectstatic → Gunicorn",
], Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5), size=14, gap=9)
tf = textbox(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.5))
setrun(tf.paragraphs[0], "Orchestration locale — 3 services", 17, PINK, bold=True)
svc = [("db", "PostgreSQL 16 + volume + healthcheck"),
       ("web", "Django servi par Gunicorn (3 workers)"),
       ("nginx", "Reverse proxy (port 80) + fichiers statiques/médias")]
cw = Inches(3.8); x0 = Inches(0.8); y = Inches(4.0)
for i, (n, d) in enumerate(svc):
    x = x0 + i * (cw + Inches(0.25))
    rect(s, x, y, cw, Inches(1.6), LIGHT, line=VIOLET)
    tf = textbox(s, x + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), Inches(1.3))
    setrun(tf.paragraphs[0], n, 16, VIOLET, bold=True, font="Consolas")
    p = tf.add_paragraph(); p.space_before = Pt(4)
    setrun(p, d, 13, DARK)
footer(s, pg())

# =============================================================== 16. CI/CD
s = slide(); header(s, "Pipeline CI/CD — GitHub Actions", "09  •  CI/CD")
box = rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.95), LIGHT, line=VIOLET)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3)
setrun(tf.paragraphs[0], "Objectif : ", 15, VIOLET, bold=True)
setrun(tf.paragraphs[0], "automatiser tests, qualité du code, build de l'image Docker "
       "et publication à chaque push sur main.", 14, DARK)
# deux jobs
jobs = [("① Tests & qualité",
         ["Checkout du code", "Python 3.11 + dépendances", "Analyse statique flake8",
          "Migrations + tests sur PostgreSQL"]),
        ("② Build & publication",
         ["Uniquement sur push → main", "Connexion à GHCR",
          "docker build + push", "Tags : latest + SHA du commit"])]
cw = Inches(5.6); x0 = Inches(0.85); y = Inches(2.75)
for i, (t, its) in enumerate(jobs):
    x = x0 + i * (cw + Inches(0.5))
    rect(s, x, y, cw, Inches(0.7), VIOLET if i == 0 else PINK)
    tf = textbox(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.7),
                 anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], t, 16, WHITE, bold=True)
    rect(s, x, y + Inches(0.7), cw, Inches(3.0), LIGHT, line=RGBColor(0xDD, 0xDD, 0xE5))
    bullets(s, its, x + Inches(0.3), y + Inches(0.95), cw - Inches(0.6), Inches(2.6),
            size=14, gap=11)
footer(s, pg())

# =============================================================== 17. SECTION DEMO
s = slide()
rect(s, 0, 0, SW, SH, VIOLET)
rect(s, 0, Inches(3.3), SW, Inches(0.9), RGBColor(0x6D, 0x28, 0xD9))
tf = textbox(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.5),
             anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
setrun(tf.paragraphs[0], "Démonstration — captures d'écran", 36, WHITE, bold=True)

# =============================================================== 18-23. SCREENSHOTS
shots = [
    ("fig-accueil-hero.png", "Page d'accueil — section « hero »",
     "Accès aux rayons Homme / Femme et mention des recommandations proposées par l'IA."),
    ("fig-accueil-homme.png", "Accueil — types & rayon Homme",
     "Types de chaussures et rayon « Chaussures Homme » sous forme de grille (prix en MAD)."),
    ("fig-catalogue.png", "Catalogue de produits",
     "Grille avec filtres (genre, types) et tri — 250 produits disponibles."),
    ("fig-panier.png", "Panier d'achat",
     "Panier (ici vide) avec invitation à découvrir les produits."),
    ("fig-dashboard.png", "Tableau de bord administrateur",
     "Indicateurs clés, meilleures ventes, commandes par statut, récentes et alertes de stock."),
    ("fig-admin.png", "Interface d'administration Django",
     "Gestion CRUD : comptes, paniers, commandes, catégories, produits et avis."),
]
for fname, title, cap in shots:
    s = slide(); header(s, title, "Démonstration")
    path = os.path.join(SS, fname)
    add_image_fit(s, path, Inches(0.7), Inches(1.45), Inches(11.95), Inches(4.85))
    tf = textbox(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.6))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    setrun(tf.paragraphs[0], cap, 13, GRAY, italic=True)
    footer(s, pg())

# =============================================================== 24. DIFFICULTES
s = slide(); header(s, "Difficultés rencontrées", "Bilan")
bullets(s, [
    "Données réelles|| : import & nettoyage du dataset Myntra (commande import_fashion)",
    "Recommandation en français|| : liste de stop-words FR constituée manuellement",
    "Performance de l'IA|| : mise en cache invalidée par la signature du catalogue",
    "Cohérence des commandes|| : transaction atomique avec vérification du stock",
    "Initialisation en conteneur|| : entrypoint idempotent (attente BDD, migrations, données)",
], Inches(0.8), Inches(1.7), Inches(11.7), Inches(5), size=16, gap=16)
footer(s, pg())

# =============================================================== 25. CONCLUSION
s = slide(); header(s, "Conclusion & perspectives", "Bilan")
tf = textbox(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(1.5))
setrun(tf.paragraphs[0], "Une plateforme e-commerce complète : ", 16, VIOLET, bold=True)
setrun(tf.paragraphs[0], "architecture modulaire, gestion des rôles, catalogue, panier, "
       "commandes, avis, tableau de bord et recommandation IA (TF-IDF + cosinus) — "
       "conteneurisée (Docker, Nginx + Gunicorn + PostgreSQL) et livrée via CI/CD.",
       16, DARK)
tf = textbox(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.5))
setrun(tf.paragraphs[0], "Perspectives", 18, PINK, bold=True)
bullets(s, [
    "Recommandation personnalisée (historique d'achat & navigation)",
    "Paiement en ligne (Stripe, PayPal)",
    "API REST (Django REST Framework) pour une application mobile",
    "Moteur de recherche avancé & gestion des promotions",
    "Déploiement continu + scans de sécurité des dépendances et de l'image",
], Inches(0.8), Inches(3.65), Inches(11.7), Inches(3), size=15, gap=11)
footer(s, pg())

# =============================================================== 26. MERCI
s = slide()
rect(s, 0, 0, SW, SH, VIOLET)
rect(s, 0, Inches(4.6), SW, Inches(0.10), PINK)
tf = textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.6),
             anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
setrun(tf.paragraphs[0], "Merci de votre attention", 44, WHITE, bold=True)
tf = textbox(s, Inches(1), Inches(4.9), Inches(11.3), Inches(1.2),
             anchor=MSO_ANCHOR.TOP)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
setrun(p, "Ali Bouhamidi  &  Said Aghrod", 20, WHITE, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
setrun(p, "Encadré par Prof. Bousselham  •  2025–2026", 15, RGBColor(0xE9, 0xD5, 0xFF))
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(14)
setrun(p, "Questions ?", 18, RGBColor(0xC4, 0xB5, 0xFD), italic=True)

out = "E-Shop_Django_Presentation.pptx"
prs.save(out)
print("OK -", out, "-", len(prs.slides._sldIdLst), "slides")
