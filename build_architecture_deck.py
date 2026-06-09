# -*- coding: utf-8 -*-
"""Genere la presentation PowerPoint AXEE ARCHITECTURE du projet E-Shop Django.
Focus : architecture du systeme, comment il est construit, et les choix techniques.
Source : rapport.tex
Sortie : E-Shop_Django_Architecture.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Charte graphique --------------------------------------------------------
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_D = RGBColor(0x5B, 0x21, 0xB6)
VIOLET_M = RGBColor(0x6D, 0x28, 0xD9)
PINK     = RGBColor(0xDB, 0x27, 0x77)
DARK     = RGBColor(0x24, 0x24, 0x2E)
GRAY     = RGBColor(0x6B, 0x6B, 0x78)
LIGHT    = RGBColor(0xF6, 0xF5, 0xFB)
CARD     = RGBColor(0xFA, 0xF9, 0xFE)
BORDER   = RGBColor(0xDD, 0xDB, 0xE8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LILAC    = RGBColor(0xE9, 0xD5, 0xFF)
LILAC_2  = RGBColor(0xC4, 0xB5, 0xFD)

# palette de categories (color-coding des couches / noeuds)
BLUE  = RGBColor(0x25, 0x63, 0xEB)
TEAL  = RGBColor(0x0D, 0x94, 0x88)
AMBER = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x15, 0x9E, 0x5B)
SLATE = RGBColor(0x47, 0x55, 0x69)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- Helpers -----------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def _no_shadow(sp):
    try:
        sp.shadow.inherit = False
    except Exception:
        pass


def rect(s, x, y, w, h, color, line=None, lw=1.0, rounded=False, radius=0.08):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    _no_shadow(sp)
    if rounded:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def chip(s, x, y, w, h, text, fill, txt=WHITE, size=11, font=FONT, bold=True):
    sp = rect(s, x, y, w, h, fill, rounded=True, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text, size, txt, bold=bold, font=font)
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def arrow(s, x, y, w, h, color, direction="right"):
    shapes = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
              "left": MSO_SHAPE.LEFT_ARROW, "up": MSO_SHAPE.UP_ARROW}
    sp = s.shapes.add_shape(shapes[direction], x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    _no_shadow(sp)
    try:
        sp.adjustments[0] = 0.55
        sp.adjustments[1] = 0.55
    except Exception:
        pass
    return sp


def connector(s, x1, y1, x2, y2, color, lw=1.6, dash=False, arrow_end=True):
    ln = s.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    _no_shadow(ln)
    lnEl = ln.line._get_or_add_ln()
    if dash:
        d = lnEl.makeelement(qn('a:prstDash'), {'val': 'dash'})
        lnEl.append(d)
    if arrow_end:
        tail = lnEl.makeelement(qn('a:tailEnd'),
                                {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnEl.append(tail)
    return ln


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.12), VIOLET)
    rect(s, 0, Inches(1.12), SW, Inches(0.055), PINK)
    # petit accent a gauche du titre
    rect(s, Inches(0.6), Inches(0.30), Inches(0.10), Inches(0.55), LILAC)
    tf = textbox(s, Inches(0.85), Inches(0.12), Inches(11.9), Inches(0.95),
                 anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        p = tf.paragraphs[0]
        run(p, kicker.upper(), 11.5, LILAC, bold=True)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    run(p2, title, 27, WHITE, bold=True)


def footer(s, page):
    tf = textbox(s, Inches(0.6), Inches(7.04), Inches(9), Inches(0.36))
    run(tf.paragraphs[0], "E-Shop Django  •  Architecture du systeme  •  2025-2026",
        9, GRAY, italic=True)
    tf2 = textbox(s, Inches(12.1), Inches(7.04), Inches(0.8), Inches(0.36))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run(tf2.paragraphs[0], str(page), 9, GRAY)


def bullets(s, items, x, y, w, h, size=16, gap=10, color=DARK, lead=PINK):
    tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = lvl
        bullet = "▸  " if lvl == 0 else "–  "
        run(p, bullet, size, lead if lvl == 0 else VIOLET, bold=True)
        if "||" in it:
            head, rest = it.split("||", 1)
            run(p, head, size, color, bold=True)
            run(p, rest, size, color)
        else:
            run(p, it, size, color)
    return tf


def node(s, x, y, w, h, title, sub=None, accent=VIOLET, fill=CARD,
         tsize=14, ssize=10.5, title_color=None, mono=False):
    """Carte 'noeud' : bandeau d'accent a gauche + titre + sous-titre."""
    rect(s, x, y, w, h, fill, line=BORDER, lw=1.0, rounded=True, radius=0.10)
    rect(s, x, y + Emu(20000), Inches(0.09), h - Emu(40000), accent, rounded=True, radius=0.5)
    tf = textbox(s, x + Inches(0.28), y, w - Inches(0.42), h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, title, tsize, title_color or accent, bold=True, font=MONO if mono else FONT)
    if sub:
        p2 = tf.add_paragraph(); p2.space_before = Pt(2)
        run(p2, sub, ssize, GRAY)
    return


PAGE = [0]
def pg():
    PAGE[0] += 1
    return PAGE[0]


# ======================================================= 1. TITRE
s = slide()
rect(s, 0, 0, SW, SH, VIOLET)
rect(s, 0, Inches(2.45), SW, Inches(2.7), VIOLET_M)
rect(s, 0, Inches(2.45), SW, Inches(0.06), PINK)
rect(s, 0, Inches(5.09), SW, Inches(0.06), PINK)

tf = textbox(s, Inches(1), Inches(0.85), Inches(11.3), Inches(0.6))
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "GENIE MIS  •  DEVELOPPEMENT WEB AVEC DJANGO",
    15, LILAC, bold=True)

tf = textbox(s, Inches(1), Inches(2.55), Inches(11.3), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "E-Shop Django", 56, WHITE, bold=True)

tf = textbox(s, Inches(1), Inches(3.75), Inches(11.3), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0],
    "Architecture du systeme — conception, construction & choix techniques",
    21, LILAC, italic=True)

tf = textbox(s, Inches(1), Inches(5.45), Inches(11.3), Inches(1.5))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Realise par  ", 17, WHITE)
run(p, "Ali Bouhamidi  &  Said Aghrod", 17, WHITE, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(7)
run(p, "Encadre par  ", 15, LILAC)
run(p, "Prof. Bousselham", 15, LILAC, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
run(p, "Annee universitaire 2025 – 2026", 13, LILAC_2)

# ======================================================= 2. AGENDA
s = slide(); header(s, "Fil conducteur", "Ce que couvre cette presentation")
items = [
    ("01", "Vue d'ensemble", "L'architecture en un coup d'oeil", VIOLET),
    ("02", "Architecture en couches", "Le patron MVT de Django", BLUE),
    ("03", "Architecture modulaire", "7 applications Django decouplees", PINK),
    ("04", "Cycle de vie d'une requete", "Du navigateur a la base de donnees", TEAL),
    ("05", "Modele de donnees", "Les entites metier & relations", AMBER),
    ("06", "Stack & justification", "Chaque choix et son pourquoi", VIOLET),
    ("07", "Brique d'IA", "Pipeline de recommandation", GREEN),
    ("08", "Securite & deploiement", "Defense en profondeur, Docker, CI/CD", SLATE),
]
cw = Inches(5.95); ch = Inches(1.18); gx = Inches(0.3); gy = Inches(0.22)
x0 = Inches(0.65); y0 = Inches(1.5)
for i, (num, t, d, col) in enumerate(items):
    r = i // 2; c = i % 2
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, CARD, line=BORDER, lw=1.0, rounded=True, radius=0.10)
    rect(s, x, y, Inches(1.0), ch, col, rounded=True, radius=0.10)
    rect(s, x + Inches(0.55), y, Inches(0.45), ch, col)
    tfn = textbox(s, x, y, Inches(1.0), ch, anchor=MSO_ANCHOR.MIDDLE)
    tfn.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tfn.paragraphs[0], num, 26, WHITE, bold=True)
    tf = textbox(s, x + Inches(1.2), y, cw - Inches(1.35), ch, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], t, 16, DARK, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(2)
    run(p, d, 11.5, GRAY)
footer(s, pg())

# ======================================================= 3. VUE D'ENSEMBLE
s = slide(); header(s, "L'architecture en un coup d'oeil", "01  •  Vue d'ensemble")
tf = textbox(s, Inches(0.7), Inches(1.32), Inches(12), Inches(0.5))
run(tf.paragraphs[0], "Une architecture web 3-tiers conteneurisee : ", 14, DARK, bold=True)
run(tf.paragraphs[0], "navigateur → reverse proxy → serveur d'application → base de donnees.",
    14, DARK)

y = Inches(2.35); h = Inches(1.55)
xs = [Inches(0.55), Inches(3.15), Inches(5.75), Inches(8.35), Inches(10.95)]
w = Inches(2.05)
nodes = [
    ("Navigateur", "Client HTTP\n(Bootstrap 5)", SLATE),
    ("Nginx", "Reverse proxy\nport 80 + statiques", TEAL),
    ("Gunicorn", "Serveur WSGI\n3 workers", AMBER),
    ("Django 5.0", "Logique metier\n(MVT)", VIOLET),
    ("PostgreSQL 16", "Donnees\n+ volume persistant", BLUE),
]
for i, (t, d, col) in enumerate(nodes):
    x = xs[i]
    rect(s, x, y, w, h, CARD, line=col, lw=1.6, rounded=True, radius=0.12)
    rect(s, x, y, w, Inches(0.52), col, rounded=True, radius=0.12)
    rect(s, x, y + Inches(0.30), w, Inches(0.22), col)
    tf = textbox(s, x, y + Inches(0.04), w, Inches(0.48), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], t, 13.5, WHITE, bold=True)
    tf = textbox(s, x + Inches(0.1), y + Inches(0.62), w - Inches(0.2), Inches(0.85),
                 anchor=MSO_ANCHOR.MIDDLE)
    for j, line in enumerate(d.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 10.5, DARK)
    if i < 4:
        ax = x + w + Emu(20000)
        arrow(s, ax, y + Inches(0.55), Inches(0.45), Inches(0.45), PINK, "right")

# legende flux retour + statiques
tf = textbox(s, Inches(0.55), y + h + Inches(0.18), Inches(12), Inches(0.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "HTTP requete →   •   ← reponse HTML / JSON   •   "
       "fichiers statiques servis par Nginx + WhiteNoise   •   medias (Pillow)",
    12, GRAY, italic=True)

# bande "tout dans Docker Compose"
by = Inches(5.55)
rect(s, Inches(0.55), by, Inches(12.25), Inches(1.05), LIGHT, line=VIOLET, lw=1.4, rounded=True, radius=0.10)
tf = textbox(s, Inches(0.9), by, Inches(11.6), Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
run(tf.paragraphs[0], "Le tout orchestre par Docker Compose", 14, VIOLET, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(2)
run(p, "3 conteneurs (nginx, web, db) sur un reseau interne  •  variables sensibles via .env  "
       "•  livraison automatisee par GitHub Actions → GHCR", 12, DARK)
footer(s, pg())

# ======================================================= 4. MVT (couches)
s = slide(); header(s, "Architecture en couches — le patron MVT", "02  •  Comment c'est construit")
tf = textbox(s, Inches(0.7), Inches(1.32), Inches(12), Inches(0.5))
run(tf.paragraphs[0], "Django impose une separation nette des responsabilites : ", 14, DARK)
run(tf.paragraphs[0], "Model – View – Template.", 14, VIOLET, bold=True)

layers = [
    ("TEMPLATE", "Presentation", "Templates Django + Bootstrap 5 — rendu HTML, responsive",
     "templates/  •  static/", BLUE),
    ("VIEW", "Logique applicative", "Vues = traitement de la requete, regles metier, orchestration",
     "views.py  •  urls.py  •  forms.py", VIOLET),
    ("MODEL", "Acces aux donnees (ORM)", "Modeles Django — mapping objet-relationnel, requetes, migrations",
     "models.py  •  migrations/", AMBER),
    ("BASE DE DONNEES", "Persistance", "PostgreSQL 16 en production (SQLite en developpement local)",
     "tables relationnelles", TEAL),
]
y = Inches(2.0); lh = Inches(1.06); gap = Inches(0.12)
for i, (tag, role, desc, code, col) in enumerate(layers):
    yy = y + i * (lh + gap)
    rect(s, Inches(0.7), yy, Inches(11.9), lh, CARD, line=BORDER, lw=1.0, rounded=True, radius=0.10)
    rect(s, Inches(0.7), yy, Inches(2.55), lh, col, rounded=True, radius=0.10)
    rect(s, Inches(2.7), yy, Inches(0.55), lh, col)
    tf = textbox(s, Inches(0.7), yy, Inches(2.55), lh, anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], tag, 14, WHITE, bold=True)
    tf = textbox(s, Inches(3.5), yy + Inches(0.08), Inches(8.9), lh - Inches(0.16),
                 anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, role + "  —  ", 13.5, DARK, bold=True)
    run(p, desc, 12.5, GRAY)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    run(p2, code, 11, col, font=MONO, bold=True)
    if i < 3:
        connector(s, Inches(6.6), yy + lh, Inches(6.6), yy + lh + gap, PINK, lw=1.6)
footer(s, pg())

# ======================================================= 5. MODULAIRE (apps)
s = slide(); header(s, "Architecture modulaire — 7 applications decouplees", "03  •  Comment c'est construit")
tf = textbox(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5))
run(tf.paragraphs[0], "Chaque domaine metier est isole dans une application Django dediee ",
    13.5, DARK)
run(tf.paragraphs[0], "→ separation des responsabilites & maintenabilite.", 13.5, VIOLET, bold=True)

# noyau config en haut
rect(s, Inches(4.4), Inches(1.95), Inches(4.5), Inches(0.7), VIOLET_M, rounded=True, radius=0.2)
tfc = textbox(s, Inches(4.4), Inches(1.95), Inches(4.5), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
tfc.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tfc.paragraphs[0], "ecommerce_project", 14, WHITE, bold=True, font=MONO)
p = tfc.add_paragraph(); p.alignment = PP_ALIGN.CENTER
run(p, "settings  •  urls  •  wsgi", 10.5, LILAC)

apps = [
    ("accounts", "Utilisateurs, inscription,\nconnexion, profils", BLUE),
    ("products", "Produits, categories,\nrecherche, filtres, avis", VIOLET),
    ("cart", "Panier d'achat\npersistant", PINK),
    ("orders", "Commandes\n& checkout", AMBER),
    ("dashboard", "Statistiques\n& administration", TEAL),
    ("recommendation", "IA — TF-IDF\n+ cosinus", GREEN),
]
cw = Inches(3.85); ch = Inches(1.5); gx = Inches(0.28); gy = Inches(0.28)
x0 = Inches(0.75); y0 = Inches(3.0)
for i, (name, desc, col) in enumerate(apps):
    r = i // 3; c = i % 3
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, CARD, line=col, lw=1.5, rounded=True, radius=0.10)
    rect(s, x, y, cw, Inches(0.16), col, rounded=True, radius=0.5)
    tf = textbox(s, x + Inches(0.25), y + Inches(0.22), cw - Inches(0.4), ch - Inches(0.3))
    run(tf.paragraphs[0], name + "/", 15, col, bold=True, font=MONO)
    for j, line in enumerate(desc.split("\n")):
        p = tf.add_paragraph()
        if j == 0:
            p.space_before = Pt(4)
        run(p, line, 11.5, DARK)
footer(s, pg())

# ======================================================= 6. CYCLE DE VIE REQUETE
s = slide(); header(s, "Cycle de vie d'une requete", "04  •  Comment ca fonctionne")
tf = textbox(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5))
run(tf.paragraphs[0], "Exemple : un client valide sa commande (checkout).", 14, DARK, italic=True)

steps = [
    ("1", "Navigateur", "Requete HTTPS\nPOST /checkout", SLATE),
    ("2", "Nginx", "Reverse proxy\n→ transmet a web", TEAL),
    ("3", "Gunicorn + URLconf", "Route l'URL vers\nla bonne vue", AMBER),
    ("4", "Middleware", "Auth, session,\nCSRF, securite", PINK),
    ("5", "Vue (View)", "@login_required\ntransaction atomique", VIOLET),
    ("6", "ORM → PostgreSQL", "Verif stock, cree\nOrder, decremente", BLUE),
    ("7", "Template", "Rend la page\nde confirmation", GREEN),
]
# 7 etapes en zig-zag 4 + 3
def step_card(x, y, num, t, d, col):
    w = Inches(2.7); h = Inches(1.5)
    rect(s, x, y, w, h, CARD, line=col, lw=1.5, rounded=True, radius=0.12)
    c = rect(s, x + Inches(0.12), y + Inches(0.12), Inches(0.5), Inches(0.5), col, rounded=True, radius=0.5)
    tfn = textbox(s, x + Inches(0.12), y + Inches(0.12), Inches(0.5), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    tfn.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tfn.paragraphs[0], num, 16, WHITE, bold=True)
    tf = textbox(s, x + Inches(0.72), y + Inches(0.12), w - Inches(0.82), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], t, 12.5, col, bold=True)
    tf = textbox(s, x + Inches(0.16), y + Inches(0.7), w - Inches(0.3), Inches(0.72))
    for j, line in enumerate(d.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        run(p, line, 10.5, DARK)
    return w, h

w = Inches(2.7); h = Inches(1.5); gx = Inches(0.33)
row1_y = Inches(2.05); row2_y = Inches(4.05)
xs1 = [Inches(0.6) + i * (w + gx) for i in range(4)]
for i in range(4):
    step_card(xs1[i], row1_y, *steps[i])
    if i < 3:
        arrow(s, xs1[i] + w + Emu(10000), row1_y + Inches(0.5), Inches(0.30), Inches(0.45), PINK, "right")
# descente
arrow(s, xs1[3] + w/2 - Inches(0.22), row1_y + h + Emu(15000), Inches(0.45), Inches(0.30), PINK, "down")
xs2 = [xs1[3] - i * (w + gx) for i in range(3)]
for i in range(3):
    step_card(xs2[i], row2_y, *steps[4 + i])
    if i < 2:
        arrow(s, xs2[i] - Inches(0.30) - Emu(10000), row2_y + Inches(0.5), Inches(0.30), Inches(0.45), PINK, "left")
footer(s, pg())

# ======================================================= 7. MODELE DE DONNEES
s = slide(); header(s, "Modele de donnees — entites & relations", "05  •  Conception")
tf = textbox(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.45))
run(tf.paragraphs[0], "Patron en-tete / lignes pour le panier et la commande ; "
    "User natif Django etendu par un Profile (1–1).", 13, DARK, italic=True)

def ent(x, y, w, h, name, fields, col):
    rect(s, x, y, w, h, CARD, line=col, lw=1.5, rounded=True, radius=0.08)
    rect(s, x, y, w, Inches(0.42), col, rounded=True, radius=0.12)
    rect(s, x, y + Inches(0.22), w, Inches(0.2), col)
    tf = textbox(s, x, y, w, Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], name, 12.5, WHITE, bold=True, font=MONO)
    tf = textbox(s, x + Inches(0.14), y + Inches(0.5), w - Inches(0.24), h - Inches(0.55))
    for j, f in enumerate(fields):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        run(p, f, 10, DARK)

ent(Inches(0.7),  Inches(2.0), Inches(2.5), Inches(1.35), "User", ["username", "email", "password (hash)", "is_staff"], SLATE)
ent(Inches(0.7),  Inches(3.7), Inches(2.5), Inches(1.2),  "Profile", ["phone, address", "city, postal_code"], BLUE)
ent(Inches(3.9),  Inches(2.0), Inches(2.5), Inches(1.5),  "Category", ["name, slug", "description", "created_at"], AMBER)
ent(Inches(3.9),  Inches(3.85),Inches(2.5), Inches(1.6),  "Product", ["name, gender, price", "image, stock", "available, slug"], VIOLET)
ent(Inches(7.1),  Inches(2.0), Inches(2.3), Inches(1.2),  "Cart", ["user (1-1)", "created_at"], PINK)
ent(Inches(7.1),  Inches(3.55),Inches(2.3), Inches(1.2),  "CartItem", ["quantity", "subtotal"], PINK)
ent(Inches(10.0), Inches(2.0), Inches(2.6), Inches(1.45), "Order", ["full_name, address", "city, phone, status"], GREEN)
ent(Inches(10.0), Inches(3.75),Inches(2.6), Inches(1.35), "OrderItem", ["product_name", "price, quantity", "subtotal"], GREEN)
ent(Inches(5.55), Inches(5.75),Inches(2.5), Inches(1.1),  "Review", ["rating 1-5, comment", "unique(user, product)"], TEAL)

# quelques relations cle
connector(s, Inches(1.95), Inches(3.7), Inches(1.95), Inches(3.35), GRAY, lw=1.4)   # profile-user
connector(s, Inches(5.15), Inches(3.85), Inches(5.15), Inches(3.5), GRAY, lw=1.4)   # product-category
connector(s, Inches(8.25), Inches(3.55), Inches(8.25), Inches(3.2), GRAY, lw=1.4)   # cartitem-cart
connector(s, Inches(11.3), Inches(3.75), Inches(11.3), Inches(3.45), GRAY, lw=1.4)  # orderitem-order
connector(s, Inches(6.4), Inches(4.6), Inches(7.1), Inches(4.15), GRAY, lw=1.4)     # cartitem-product
tf = textbox(s, Inches(0.7), Inches(6.95), Inches(8), Inches(0.4))
run(tf.paragraphs[0], "Prix et nom figes dans OrderItem a l'instant de l'achat.", 11, GRAY, italic=True)
footer(s, pg())

# ======================================================= 8. STACK & JUSTIFICATION
s = slide(); header(s, "Stack technologique & justification des choix", "06  •  Les choix")
rows = [
    ("Backend", "Django 5.0 / Python 3.11", "Framework batteries-included, ORM, admin & securite integres", VIOLET),
    ("Frontend", "Templates + Bootstrap 5", "Rendu serveur simple, responsive, pas de SPA a maintenir", BLUE),
    ("Base de donnees", "PostgreSQL 16 (SQLite en dev)", "SGBD robuste pour la prod, demarrage instantane en local", TEAL),
    ("Fichiers statiques", "WhiteNoise (manifest + gzip)", "Sert les statiques sans service externe, cache long terme", AMBER),
    ("Intelligence artif.", "scikit-learn (TF-IDF + cosinus)", "Recommandation content-based legere, sans GPU ni API tierce", GREEN),
    ("Serveur d'app", "Gunicorn (3 workers WSGI)", "Serveur WSGI eprouve, concurrence par workers", AMBER),
    ("Reverse proxy", "Nginx", "Sert statiques/medias, point d'entree unique, TLS-ready", TEAL),
    ("Conteneurisation", "Docker & Docker Compose", "Environnement reproductible, parite dev / prod", BLUE),
    ("CI/CD & registre", "GitHub Actions → GHCR", "Tests + build + publication automatises a chaque push", PINK),
]
y0 = Inches(1.4); rh = Inches(0.55); x0 = Inches(0.6)
w1 = Inches(2.65); w2 = Inches(3.55); w3 = Inches(5.95)
# en-tete
rect(s, x0, y0, w1, rh, VIOLET_M); rect(s, x0 + w1, y0, w2, rh, VIOLET_M); rect(s, x0 + w1 + w2, y0, w3, rh, VIOLET_M)
for lbl, xx, ww in [("Couche", x0, w1), ("Technologie retenue", x0 + w1, w2), ("Pourquoi ce choix", x0 + w1 + w2, w3)]:
    tf = textbox(s, xx + Inches(0.15), y0, ww - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], lbl, 12.5, WHITE, bold=True)
for i, (a, b, c, col) in enumerate(rows):
    y = y0 + (i + 1) * rh
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, x0, y, w1 + w2 + w3, rh, bg, line=BORDER, lw=0.75)
    rect(s, x0, y, Inches(0.07), rh, col)
    tf = textbox(s, x0 + Inches(0.2), y, w1 - Inches(0.25), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], a, 11.5, DARK, bold=True)
    tf = textbox(s, x0 + w1 + Inches(0.12), y, w2 - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], b, 11.5, col, bold=True)
    tf = textbox(s, x0 + w1 + w2 + Inches(0.12), y, w3 - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], c, 11, GRAY)
footer(s, pg())

# ======================================================= 9. PIPELINE IA
s = slide(); header(s, "Brique d'IA — pipeline de recommandation", "07  •  Comment c'est construit")
box = rect(s, Inches(0.7), Inches(1.32), Inches(11.9), Inches(0.85), LIGHT, line=VIOLET, lw=1.3, rounded=True, radius=0.08)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.25)
run(tf.paragraphs[0], "Recommandation de produits similaires ", 14, VIOLET, bold=True)
run(tf.paragraphs[0], "(content-based filtering) — bloc « Vous aimerez aussi » sur chaque fiche produit.",
    13, DARK)

pipe = [
    ("Produit → document", "nom (pondere) +\ncategorie + description", BLUE),
    ("TF-IDF (fit)", "entraine sur 20% du\ncatalogue, n-grammes 1-2", VIOLET),
    ("transform", "projection du catalogue\ncomplet dans l'espace", AMBER),
    ("Similarite cosinus", "matrice produit-produit\n→ les k plus proches", GREEN),
]
w = Inches(2.7); h = Inches(1.7); gx = Inches(0.37); y = Inches(2.55)
xs = [Inches(0.7) + i * (w + gx) for i in range(4)]
for i, (t, d, col) in enumerate(pipe):
    x = xs[i]
    rect(s, x, y, w, h, CARD, line=col, lw=1.6, rounded=True, radius=0.10)
    rect(s, x, y, w, Inches(0.5), col, rounded=True, radius=0.12)
    rect(s, x, y + Inches(0.28), w, Inches(0.22), col)
    tf = textbox(s, x, y, w, Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], t, 12.5, WHITE, bold=True, font=MONO)
    tf = textbox(s, x + Inches(0.12), y + Inches(0.62), w - Inches(0.24), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    for j, line in enumerate(d.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 11, DARK)
    if i < 3:
        arrow(s, x + w + Emu(15000), y + Inches(0.62), Inches(0.32), Inches(0.45), PINK, "right")

# optimisation
oy = Inches(4.75)
rect(s, Inches(0.7), oy, Inches(5.85), Inches(1.55), LIGHT, line=TEAL, lw=1.3, rounded=True, radius=0.08)
tf = textbox(s, Inches(0.95), oy + Inches(0.12), Inches(5.4), Inches(1.3))
run(tf.paragraphs[0], "Optimisation — mise en cache", 13, TEAL, bold=True)
bullets(s, ["Modele reconstruit seulement si la signature du catalogue change",
            "Evite de recalculer la matrice a chaque requete"],
        Inches(0.95), oy + Inches(0.5), Inches(5.4), Inches(1.0), size=11, gap=5, lead=TEAL)

rect(s, Inches(6.75), oy, Inches(5.85), Inches(1.55), LIGHT, line=AMBER, lw=1.3, rounded=True, radius=0.08)
tf = textbox(s, Inches(7.0), oy + Inches(0.12), Inches(5.4), Inches(1.3))
run(tf.paragraphs[0], "Robustesse — repli (fallback)", 13, AMBER, bold=True)
bullets(s, ["Aucun resultat ? → repli sur les produits de la meme categorie",
            "Stop-words FR constitues manuellement (absents de scikit-learn)"],
        Inches(7.0), oy + Inches(0.5), Inches(5.4), Inches(1.0), size=11, gap=5, lead=AMBER)
footer(s, pg())

# ======================================================= 10. SECURITE
s = slide(); header(s, "Architecture de securite — defense en profondeur", "08  •  Les choix")
cards = [
    ("Authentification & roles", BLUE,
     ["@login_required : panier, checkout, avis",
      "is_staff : dashboard & espace admin",
      "Un client n'accede qu'a ses commandes"]),
    ("Protection des donnees", VIOLET,
     ["Mots de passe haches (+ 4 validators)",
      "Protection CSRF + validation serveur",
      "SECRET_KEY & BDD dans .env non versionne"]),
    ("Durcissement production", PINK,
     ["DEBUG=False, cookies securises",
      "HSTS, X-Frame-Options: DENY, nosniff",
      "Uploads limites : 2,5 Mo / 5 Mo (Nginx)"]),
]
cw = Inches(3.9); gap = Inches(0.3); x0 = Inches(0.65); y0 = Inches(1.55); ch = Inches(4.4)
for i, (t, col, its) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, CARD, line=col, lw=1.5, rounded=True, radius=0.07)
    rect(s, x, y0, cw, Inches(0.95), col, rounded=True, radius=0.10)
    rect(s, x, y0 + Inches(0.6), cw, Inches(0.35), col)
    tf = textbox(s, x + Inches(0.2), y0, cw - Inches(0.4), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], t, 15, WHITE, bold=True)
    bullets(s, its, x + Inches(0.28), y0 + Inches(1.2), cw - Inches(0.5), Inches(3.0),
            size=12.5, gap=14, lead=col)
tf = textbox(s, Inches(0.65), Inches(6.25), Inches(12), Inches(0.6))
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "Les mots de passe ne sont jamais stockes en clair  •  le .env n'est jamais publie (.gitignore)",
    12, VIOLET, bold=True, italic=True)
footer(s, pg())

# ======================================================= 11. DEPLOIEMENT
s = slide(); header(s, "Conteneurisation & deploiement", "08  •  Comment c'est livre")
tf = textbox(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.45))
run(tf.paragraphs[0], "Image basee sur python:3.11-slim. ", 13.5, DARK, bold=True)
run(tf.paragraphs[0], "Docker Compose orchestre 3 services sur un reseau interne.", 13.5, DARK)

svc = [
    ("nginx", "Reverse proxy", ["Expose le port 80", "Sert statiques & medias", "→ transmet a web"], TEAL),
    ("web", "Django + Gunicorn", ["3 workers WSGI", "entrypoint.sh idempotent", "logique applicative"], VIOLET),
    ("db", "PostgreSQL 16", ["Volume persistant", "Healthcheck", "reseau interne uniquement"], BLUE),
]
cw = Inches(3.55); gap = Inches(0.35); x0 = Inches(0.95); y0 = Inches(2.0); ch = Inches(2.6)
for i, (n, role, its, col) in enumerate(svc):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, CARD, line=col, lw=1.6, rounded=True, radius=0.08)
    rect(s, x, y0, cw, Inches(0.75), col, rounded=True, radius=0.12)
    rect(s, x, y0 + Inches(0.5), cw, Inches(0.25), col)
    tf = textbox(s, x, y0 + Inches(0.06), cw, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], n, 16, WHITE, bold=True, font=MONO)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    run(p, role, 10.5, LILAC)
    bullets(s, its, x + Inches(0.25), y0 + Inches(0.95), cw - Inches(0.45), Inches(1.5),
            size=11.5, gap=8, lead=col)
    if i < 2:
        arrow(s, x + cw + Emu(10000), y0 + Inches(1.05), Inches(0.32), Inches(0.5), PINK, "right")

# entrypoint flow
ey = Inches(5.05)
rect(s, Inches(0.7), ey, Inches(11.9), Inches(1.45), LIGHT, line=AMBER, lw=1.3, rounded=True, radius=0.08)
tf = textbox(s, Inches(0.95), ey + Inches(0.1), Inches(11.4), Inches(0.4))
run(tf.paragraphs[0], "entrypoint.sh (sequence de demarrage idempotente)", 13, AMBER, bold=True)
seq = ["Attendre\nPostgreSQL", "Appliquer\nles migrations", "Charger les\ndonnees (si vide)", "collectstatic", "Lancer\nGunicorn"]
sw = Inches(1.95); sx = Inches(0.95); sy = ey + Inches(0.55)
for i, st in enumerate(seq):
    x = sx + i * (sw + Inches(0.35))
    rect(s, x, sy, sw, Inches(0.75), WHITE, line=AMBER, lw=1.2, rounded=True, radius=0.12)
    tf = textbox(s, x + Inches(0.05), sy, sw - Inches(0.1), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    for j, line in enumerate(st.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 10.5, DARK, bold=(j == 0))
    if i < 4:
        arrow(s, x + sw + Emu(8000), sy + Inches(0.22), Inches(0.27), Inches(0.32), AMBER, "right")
footer(s, pg())

# ======================================================= 12. CI/CD
s = slide(); header(s, "Pipeline CI/CD — GitHub Actions", "08  •  Comment c'est livre")
box = rect(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.8), LIGHT, line=VIOLET, lw=1.3, rounded=True, radius=0.08)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.25)
run(tf.paragraphs[0], "Objectif : ", 14, VIOLET, bold=True)
run(tf.paragraphs[0], "automatiser tests, qualite, build de l'image Docker et publication a chaque push sur main.",
    13, DARK)

# trigger
rect(s, Inches(0.7), Inches(2.6), Inches(2.3), Inches(3.0), SLATE, rounded=True, radius=0.06)
tf = textbox(s, Inches(0.7), Inches(2.6), Inches(2.3), Inches(3.0), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "git push", 16, WHITE, bold=True, font=MONO)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(4)
run(p, "→ main", 13, LILAC, font=MONO)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(12)
run(p, "declenche\nle workflow", 11, LILAC_2)

arrow(s, Inches(3.05), Inches(3.85), Inches(0.45), Inches(0.55), PINK, "right")

jobs = [
    ("JOB 1  —  Tests & qualite", VIOLET,
     ["Checkout du code", "Python 3.11 + dependances", "Analyse statique flake8",
      "Migrations + tests (service PostgreSQL)"]),
    ("JOB 2  —  Build & publication", PINK,
     ["Uniquement sur push → main  (needs: test)", "Connexion a GHCR",
      "docker build + push", "Tags : latest + SHA du commit"]),
]
jx = Inches(3.65); jw = Inches(4.35); jy = Inches(2.6); jh = Inches(3.0)
for i, (t, col, its) in enumerate(jobs):
    x = jx + i * (jw + Inches(0.25))
    rect(s, x, jy, jw, jh, CARD, line=col, lw=1.6, rounded=True, radius=0.06)
    rect(s, x, jy, jw, Inches(0.62), col, rounded=True, radius=0.10)
    rect(s, x, jy + Inches(0.4), jw, Inches(0.22), col)
    tf = textbox(s, x, jy, jw, Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf.paragraphs[0], t, 14, WHITE, bold=True)
    bullets(s, its, x + Inches(0.28), jy + Inches(0.82), jw - Inches(0.5), Inches(2.1),
            size=12, gap=12, lead=col)
    if i == 0:
        arrow(s, x + jw + Emu(5000), jy + Inches(1.1), Inches(0.32), Inches(0.55), PINK, "right")
tf = textbox(s, Inches(3.65), Inches(5.8), Inches(9), Inches(0.45))
run(tf.paragraphs[0], "Resultat : une image Docker testee & publiee, prete a deployer.",
    12, GRAY, italic=True)
footer(s, pg())

# ======================================================= 13. DECISIONS / RECAP
s = slide(); header(s, "Decisions d'architecture — synthese", "Recapitulatif des choix")
dec = [
    ("Architecture modulaire (apps Django)", "Separation des responsabilites, code reutilisable et testable", VIOLET),
    ("Patron MVT strict", "Logique metier dans les vues, jamais dans les templates", BLUE),
    ("Panier persistant en base", "Le panier survit a la deconnexion (1 Cart / utilisateur)", PINK),
    ("Checkout en transaction atomique", "Coherence garantie : verif stock + commande + decrement", AMBER),
    ("Prix & nom figes dans OrderItem", "L'historique reste exact meme si le produit change", GREEN),
    ("Cache du modele IA par signature", "Performance : pas de recalcul inutile de la similarite", TEAL),
    ("Repli par categorie (IA)", "Une recommandation pertinente en toute circonstance", VIOLET),
    ("Conteneurisation + CI/CD", "Parite dev/prod, livraison automatisee et reproductible", SLATE),
]
cw = Inches(5.95); ch = Inches(1.18); gx = Inches(0.3); gy = Inches(0.16)
x0 = Inches(0.65); y0 = Inches(1.45)
for i, (t, d, col) in enumerate(dec):
    r = i // 2; c = i % 2
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, CARD, line=BORDER, lw=1.0, rounded=True, radius=0.10)
    rect(s, x, y, Inches(0.1), ch, col, rounded=True, radius=0.5)
    tf = textbox(s, x + Inches(0.3), y + Inches(0.12), cw - Inches(0.5), ch - Inches(0.2),
                 anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], t, 13.5, col, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(2)
    run(p, d, 11.5, GRAY)
footer(s, pg())

# ======================================================= 14. MERCI
s = slide()
rect(s, 0, 0, SW, SH, VIOLET)
rect(s, 0, Inches(2.4), SW, Inches(2.3), VIOLET_M)
rect(s, 0, Inches(2.4), SW, Inches(0.06), PINK)
rect(s, 0, Inches(4.64), SW, Inches(0.06), PINK)
tf = textbox(s, Inches(1), Inches(2.55), Inches(11.3), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "Merci de votre attention", 44, WHITE, bold=True)
tf = textbox(s, Inches(1), Inches(3.65), Inches(11.3), Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0], "Une architecture modulaire, securisee, conteneurisee et livree en continu.",
    17, LILAC, italic=True)
tf = textbox(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Ali Bouhamidi  &  Said Aghrod", 20, WHITE, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
run(p, "Encadre par Prof. Bousselham  •  2025–2026", 14, LILAC)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(14)
run(p, "Questions ?", 17, LILAC_2, italic=True)

out = "E-Shop_Django_Architecture.pptx"
prs.save(out)
print("OK -", out, "-", len(prs.slides._sldIdLst), "slides")
